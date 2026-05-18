from __future__ import annotations

import json
from typing import Any

import httpx
import pandas as pd
from ddgs import DDGS
from docx import Document
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches as DocxInches
from docx.shared import Pt as DocxPt
from docx.shared import RGBColor
from langchain_core.tools import tool
from openpyxl import Workbook
from pypdf import PdfReader
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx import Presentation
from pptx.util import Inches as PptxInches
from pptx.util import Pt as PptxPt

from .memory import MemoryStore
from .runlog import append_log
from .sandbox import Sandbox, slugify


def build_tools(sandbox: Sandbox, memory: MemoryStore):
    @tool
    def think(title: str, thought: str, action: str = "", confidence: float = 0.8) -> str:
        """Record a concise progress plan before tool work. Do not include private chain-of-thought."""
        append_log("reasoning", title, f"Next: {action or 'continue'} · confidence {confidence:.2f}")
        return f"Recorded plan: {title}. Next action: {action or 'continue'}."

    @tool
    def analyze(title: str, result: str, analysis: str, next_action: str = "continue", confidence: float = 0.8) -> str:
        """Record a concise analysis checkpoint after tool results. Do not include private chain-of-thought."""
        append_log(
            "analysis",
            title,
            f"Result: {_short(result, 220)} | Next: {next_action} · confidence {confidence:.2f}",
        )
        return f"Recorded analysis: {title}. Next action: {next_action}."

    @tool
    def list_uploaded_files() -> str:
        """List files uploaded into the sandbox."""
        append_log("tool", "Listing uploaded files", "Scanning sandbox uploads.")
        sandbox.ensure()
        files = []
        for path in sorted(sandbox.uploads_dir.rglob("*")):
            if path.is_file():
                files.append(
                    {
                        "path": str(path.relative_to(sandbox.uploads_dir)),
                        "size": path.stat().st_size,
                    }
                )
        return json.dumps(files, indent=2)

    @tool
    def read_text_file(path: str, max_chars: int = 12000) -> str:
        """Read a text, PDF, DOCX, PPTX, XLSX, CSV, or markdown file from the sandbox."""
        append_log("tool", "Reading file", path)
        file_path = sandbox.resolve_read(path)
        if not file_path.exists():
            return f"error: file not found: {path}"
        suffix = file_path.suffix.lower()
        try:
            if suffix == ".pdf":
                reader = PdfReader(str(file_path))
                text = "\n\n".join(page.extract_text() or "" for page in reader.pages)
            elif suffix == ".docx":
                doc = Document(str(file_path))
                text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            elif suffix == ".pptx":
                prs = Presentation(str(file_path))
                slides = []
                for index, slide in enumerate(prs.slides, start=1):
                    parts = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            parts.append(shape.text.strip())
                    slides.append(f"Slide {index}\n" + "\n".join(parts))
                text = "\n\n".join(slides)
            elif suffix in {".xlsx", ".xlsm"}:
                sheets = pd.read_excel(file_path, sheet_name=None, nrows=20)
                chunks = []
                for name, frame in sheets.items():
                    chunks.append(f"Sheet: {name}\n{frame.to_csv(index=False)}")
                text = "\n\n".join(chunks)
            elif suffix == ".csv":
                frame = pd.read_csv(file_path, nrows=50)
                text = frame.to_csv(index=False)
            else:
                text = file_path.read_text(encoding="utf-8", errors="ignore")
            return text[:max_chars]
        except Exception as exc:
            return f"error reading {path}: {exc}"

    @tool
    def write_markdown(filename: str, content: str) -> str:
        """Write a markdown or text artifact under sandbox outputs."""
        append_log("artifact", "Writing markdown", filename)
        safe_name = slugify(filename, "artifact.md")
        if not safe_name.endswith((".md", ".txt")):
            safe_name += ".md"
        target = sandbox.resolve_output(safe_name)
        target.write_text(content.rstrip() + "\n", encoding="utf-8")
        return json.dumps({"status": "success", "path": str(target)})

    @tool
    def duckduckgo_search(query: str, max_results: int = 6) -> str:
        """Search the web using DuckDuckGo via ddgs. Use for research discovery."""
        append_log("research", "Searching DuckDuckGo", query)
        max_results = max(1, min(max_results, 10))
        try:
            with DDGS() as ddgs:
                results = list(ddgs.text(query, max_results=max_results))
            normalized = [
                {
                    "title": item.get("title"),
                    "href": item.get("href") or item.get("url"),
                    "body": item.get("body"),
                }
                for item in results
            ]
            return json.dumps(normalized, indent=2)
        except Exception as exc:
            return f"error: DuckDuckGo search failed: {exc}"

    @tool
    def fetch_url(url: str, max_chars: int = 14000) -> str:
        """Fetch a URL and return text content. Use after search to inspect sources."""
        append_log("research", "Fetching URL", url)
        try:
            headers = {"User-Agent": "code-assist-langgraph-prototype/0.1"}
            with httpx.Client(follow_redirects=True, timeout=20, headers=headers) as client:
                response = client.get(url)
                response.raise_for_status()
            content_type = response.headers.get("content-type", "")
            text = response.text
            return json.dumps(
                {
                    "url": str(response.url),
                    "status_code": response.status_code,
                    "content_type": content_type,
                    "text": text[:max_chars],
                },
                indent=2,
            )
        except Exception as exc:
            return f"error: fetch failed for {url}: {exc}"

    @tool
    def create_word_doc(filename: str, title: str, sections_json: str) -> str:
        """Create a Word DOCX. sections_json must be a JSON list of {heading, paragraphs}."""
        append_log("artifact", "Creating Word document", filename)
        safe_name = slugify(filename, "document")
        if not safe_name.endswith(".docx"):
            safe_name += ".docx"
        target = sandbox.resolve_output(safe_name)
        sections = _load_sections(sections_json)
        doc = Document()
        _apply_report_style(doc)
        _add_accent_title(doc, title, subtitle="")
        for section in sections:
            heading = section.get("heading")
            if heading:
                _add_section_heading(doc, str(heading))
            paragraphs = section.get("paragraphs", [])
            if isinstance(paragraphs, str):
                paragraphs = [paragraphs]
            for paragraph in paragraphs:
                _add_body_paragraph(doc, str(paragraph))
        doc.save(target)
        return json.dumps({"status": "success", "path": str(target)})

    @tool
    def create_analyst_word_report(filename: str, report_json: str) -> str:
        """Create a polished English analyst-style Word report.

        report_json is a JSON object with optional keys:
        kicker, title, ticker, subtitle, price, price_note, date, metrics,
        thesis_title, thesis_paragraphs, callout, sections, tables,
        risk_cards, conclusion, disclaimer, footer.
        """
        append_log("artifact", "Creating analyst Word report", filename)
        safe_name = slugify(filename, "analyst-report")
        if not safe_name.endswith(".docx"):
            safe_name += ".docx"
        target = sandbox.resolve_output(safe_name)
        try:
            report = json.loads(report_json)
        except json.JSONDecodeError as exc:
            raise ValueError(f"invalid report_json: {exc}") from exc
        if not isinstance(report, dict):
            raise ValueError("report_json must be a JSON object")

        doc = Document()
        _apply_report_style(doc)
        _add_report_header(doc, report)
        _add_metrics_strip(doc, report.get("metrics", []))

        thesis_title = str(report.get("thesis_title") or "Investment Thesis")
        _add_section_heading(doc, thesis_title)
        for paragraph in _as_list(report.get("thesis_paragraphs", [])):
            _add_body_paragraph(doc, paragraph)
        if report.get("callout"):
            _add_callout(doc, str(report["callout"]))

        for section in report.get("sections", []) or []:
            if not isinstance(section, dict):
                continue
            _add_section_heading(doc, str(section.get("heading", "Section")))
            for paragraph in _as_list(section.get("paragraphs", [])):
                _add_body_paragraph(doc, paragraph)
            for subsection in section.get("subsections", []) or []:
                if not isinstance(subsection, dict):
                    continue
                _add_subheading(doc, str(subsection.get("heading", "Subsection")))
                for paragraph in _as_list(subsection.get("paragraphs", [])):
                    _add_body_paragraph(doc, paragraph)

        for table_spec in report.get("tables", []) or []:
            if isinstance(table_spec, dict):
                _add_data_table(doc, table_spec)

        risk_cards = report.get("risk_cards", []) or []
        if risk_cards:
            _add_section_heading(doc, "Risks And Summary")
            _add_risk_cards(doc, risk_cards)

        if report.get("conclusion"):
            _add_callout(doc, str(report["conclusion"]), title="Conclusion")

        _add_footer_note(
            doc,
            disclaimer=str(report.get("disclaimer") or "For informational use only. Not investment advice."),
            footer=str(report.get("footer") or ""),
        )
        doc.save(target)
        return json.dumps({"status": "success", "path": str(target)})

    @tool
    def create_powerpoint(filename: str, title: str, slides_json: str) -> str:
        """Create a warm editorial PPTX deck. slides_json is a JSON list of slide objects."""
        append_log("artifact", "Creating PowerPoint deck", filename)
        safe_name = slugify(filename, "presentation")
        if not safe_name.endswith(".pptx"):
            safe_name += ".pptx"
        target = sandbox.resolve_output(safe_name)
        slides = _load_sections(slides_json)
        _create_warm_editorial_deck(title=title, slides=slides, target=target)
        return json.dumps({"status": "success", "path": str(target), "slides": len(slides) + 1})

    @tool
    def create_excel_workbook(filename: str, sheets_json: str) -> str:
        """Create an XLSX workbook. sheets_json must be JSON list of {name, rows}; rows is a list of objects."""
        append_log("artifact", "Creating Excel workbook", filename)
        safe_name = slugify(filename, "workbook")
        if not safe_name.endswith(".xlsx"):
            safe_name += ".xlsx"
        target = sandbox.resolve_output(safe_name)
        sheets = _load_sections(sheets_json)
        wb = Workbook()
        default = wb.active
        wb.remove(default)
        for sheet in sheets:
            name = str(sheet.get("name", "Sheet"))[:31]
            rows = sheet.get("rows", [])
            ws = wb.create_sheet(title=name or "Sheet")
            if rows and isinstance(rows, list) and isinstance(rows[0], dict):
                headers = list(rows[0].keys())
                ws.append(headers)
                for row in rows:
                    ws.append([row.get(header) for header in headers])
            elif isinstance(rows, list):
                for row in rows:
                    ws.append(row if isinstance(row, list) else [row])
        wb.save(target)
        return json.dumps({"status": "success", "path": str(target)})

    @tool
    def run_ppt_skill(markdown_outline: str, output_name: str = "presentation") -> str:
        """Use the local PPT skill to create a PPTX from a markdown outline."""
        append_log("artifact", "Running PPT skill", output_name)
        safe_name = slugify(output_name, "presentation")
        outline_path = sandbox.resolve_scratch(f"{safe_name}-outline.md")
        deck_path = sandbox.resolve_output(f"{safe_name}.pptx")
        outline_path.write_text(markdown_outline.rstrip() + "\n", encoding="utf-8")
        try:
            title, slides = _parse_markdown_slides(markdown_outline)
            create_powerpoint.invoke(
                {
                    "filename": deck_path.name,
                    "title": title,
                    "slides_json": json.dumps(slides),
                }
            )
        except Exception as exc:
            return json.dumps({"status": "error", "error": str(exc)})
        return json.dumps(
            {
                "status": "success",
                "outline": str(outline_path),
                "path": str(deck_path),
                "slides": len(slides) + 1,
            },
            indent=2,
        )

    @tool
    def remember(namespace: str, content: str) -> str:
        """Persist a user/project memory. Use only when the user asks to remember something or the fact is clearly durable."""
        append_log("memory", "Writing memory", namespace)
        memory_id = memory.remember(namespace=namespace or "default", content=content)
        return json.dumps({"status": "success", "id": memory_id})

    @tool
    def search_memory(namespace: str = "default", query: str = "", limit: int = 10) -> str:
        """Search persisted memory for a namespace."""
        append_log("memory", "Searching memory", f"{namespace}: {query}")
        return json.dumps(memory.search(namespace=namespace, query=query, limit=limit), indent=2)

    return [
        think,
        analyze,
        list_uploaded_files,
        read_text_file,
        write_markdown,
        duckduckgo_search,
        fetch_url,
        create_word_doc,
        create_analyst_word_report,
        create_powerpoint,
        create_excel_workbook,
        run_ppt_skill,
        remember,
        search_memory,
    ]


def _short(value: str, limit: int) -> str:
    text = str(value or "")
    if len(text) <= limit:
        return text
    return text[: limit - 1].rstrip() + "..."


def _load_sections(raw: str) -> list[dict[str, Any]]:
    try:
        parsed = json.loads(raw)
    except json.JSONDecodeError as exc:
        raise ValueError(f"invalid JSON: {exc}") from exc
    if not isinstance(parsed, list):
        raise ValueError("JSON must be a list")
    return [item if isinstance(item, dict) else {"value": item} for item in parsed]


def _as_list(value: Any) -> list[str]:
    if value is None:
        return []
    if isinstance(value, str):
        return [value]
    if isinstance(value, list):
        return [str(item) for item in value]
    return [str(value)]


def _apply_report_style(doc: Document) -> None:
    section = doc.sections[0]
    section.top_margin = DocxInches(0.65)
    section.bottom_margin = DocxInches(0.55)
    section.left_margin = DocxInches(0.72)
    section.right_margin = DocxInches(0.72)

    styles = doc.styles
    styles["Normal"].font.name = "Aptos"
    styles["Normal"].font.size = DocxPt(10.5)
    styles["Normal"].font.color.rgb = RGBColor(25, 25, 25)


def _add_report_header(doc: Document, report: dict[str, Any]) -> None:
    table = doc.add_table(rows=1, cols=2)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.autofit = True
    left, right = table.rows[0].cells
    _clear_cell(left)
    _clear_cell(right)
    _set_cell_border(left, left={"val": "single", "sz": "18", "color": "183A63"})

    kicker = str(report.get("kicker") or "Analyst Brief")
    title = str(report.get("title") or "Company Report")
    ticker = str(report.get("ticker") or "")
    subtitle = str(report.get("subtitle") or "")

    p = left.paragraphs[0]
    p.paragraph_format.left_indent = DocxInches(0.12)
    run = p.add_run(kicker.upper())
    run.font.size = DocxPt(8.5)
    run.font.bold = True
    run.font.color.rgb = RGBColor(24, 58, 99)
    p.add_run("\n")
    title_run = p.add_run(title)
    title_run.font.size = DocxPt(26)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(20, 20, 20)
    if ticker:
        ticker_run = p.add_run(f"  {ticker}")
        ticker_run.font.size = DocxPt(13)
        ticker_run.font.bold = True
        ticker_run.font.color.rgb = RGBColor(120, 120, 112)
    if subtitle:
        p.add_run("\n")
        subtitle_run = p.add_run(subtitle)
        subtitle_run.font.size = DocxPt(10.5)
        subtitle_run.font.color.rgb = RGBColor(92, 92, 86)

    price = str(report.get("price") or "")
    price_note = str(report.get("price_note") or "")
    date = str(report.get("date") or "")
    p = right.paragraphs[0]
    p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    if price:
        run = p.add_run(price)
        run.font.size = DocxPt(27)
        run.font.bold = True
    if price_note:
        p.add_run("\n")
        note_run = p.add_run(price_note)
        note_run.font.size = DocxPt(10)
        note_run.font.bold = True
        note_run.font.color.rgb = RGBColor(102, 102, 96)
    if date:
        p.add_run("\n")
        date_run = p.add_run(date)
        date_run.font.size = DocxPt(9.5)
        date_run.font.color.rgb = RGBColor(120, 120, 112)


def _add_accent_title(doc: Document, title: str, subtitle: str = "") -> None:
    _add_report_header(doc, {"title": title, "subtitle": subtitle, "kicker": "Document"})


def _add_metrics_strip(doc: Document, metrics: Any) -> None:
    if not isinstance(metrics, list) or not metrics:
        return
    cols = min(max(len(metrics), 1), 4)
    table = doc.add_table(rows=1, cols=cols)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.autofit = True
    for index, metric in enumerate(metrics[:cols]):
        cell = table.rows[0].cells[index]
        _clear_cell(cell)
        p = cell.paragraphs[0]
        value = str(metric.get("value", "")) if isinstance(metric, dict) else str(metric)
        label = str(metric.get("label", "")) if isinstance(metric, dict) else ""
        note = str(metric.get("note", "")) if isinstance(metric, dict) else ""
        run = p.add_run(value)
        run.font.size = DocxPt(17)
        run.font.bold = True
        run.font.color.rgb = RGBColor(24, 58, 99)
        if label:
            label_run = p.add_run(f"  {label}")
            label_run.font.size = DocxPt(9)
            label_run.font.color.rgb = RGBColor(92, 92, 86)
        if note:
            p.add_run("\n")
            note_run = p.add_run(note)
            note_run.font.size = DocxPt(8.5)
            note_run.font.color.rgb = RGBColor(120, 120, 112)
    doc.add_paragraph()


def _add_section_heading(doc: Document, text: str) -> None:
    p = doc.add_paragraph()
    p.paragraph_format.space_before = DocxPt(12)
    p.paragraph_format.space_after = DocxPt(4)
    run = p.add_run(text)
    run.font.size = DocxPt(17)
    run.font.bold = True
    run.font.color.rgb = RGBColor(20, 20, 20)


def _add_subheading(doc: Document, text: str) -> None:
    p = doc.add_paragraph()
    p.paragraph_format.space_before = DocxPt(8)
    p.paragraph_format.space_after = DocxPt(1)
    run = p.add_run(text)
    run.font.size = DocxPt(12.5)
    run.font.bold = True
    run.font.color.rgb = RGBColor(55, 55, 52)


def _add_body_paragraph(doc: Document, text: str) -> None:
    p = doc.add_paragraph()
    p.paragraph_format.line_spacing = 1.15
    p.paragraph_format.space_after = DocxPt(5)
    run = p.add_run(text)
    run.font.size = DocxPt(10.5)


def _add_callout(doc: Document, text: str, title: str | None = None) -> None:
    table = doc.add_table(rows=1, cols=1)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    cell = table.rows[0].cells[0]
    _shade_cell(cell, "FBFAF6")
    _set_cell_border(cell, left={"val": "single", "sz": "18", "color": "183A63"})
    p = cell.paragraphs[0]
    p.paragraph_format.left_indent = DocxInches(0.1)
    if title:
        run = p.add_run(title)
        run.font.bold = True
        run.font.color.rgb = RGBColor(24, 58, 99)
        p.add_run("\n")
    body = p.add_run(text)
    body.font.size = DocxPt(10.5)
    doc.add_paragraph()


def _add_data_table(doc: Document, table_spec: dict[str, Any]) -> None:
    title = table_spec.get("title")
    if title:
        _add_section_heading(doc, str(title))
    columns = [str(col) for col in table_spec.get("columns", [])]
    rows = table_spec.get("rows", [])
    if not columns and rows and isinstance(rows[0], dict):
        columns = [str(key) for key in rows[0].keys()]
    if not columns:
        return

    table = doc.add_table(rows=1, cols=len(columns))
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = "Table Grid"
    header = table.rows[0].cells
    for index, column in enumerate(columns):
        _shade_cell(header[index], "F2F0EA")
        header[index].vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
        p = header[index].paragraphs[0]
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER if index else WD_ALIGN_PARAGRAPH.LEFT
        run = p.add_run(column)
        run.font.bold = True
        run.font.size = DocxPt(9.5)
        run.font.color.rgb = RGBColor(75, 75, 70)

    for row in rows:
        cells = table.add_row().cells
        values = [row.get(col, "") for col in columns] if isinstance(row, dict) else row
        for index, value in enumerate(values[: len(columns)]):
            p = cells[index].paragraphs[0]
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER if index else WD_ALIGN_PARAGRAPH.LEFT
            run = p.add_run(str(value))
            run.font.size = DocxPt(9.5)
    if table_spec.get("note"):
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run(str(table_spec["note"]))
        run.font.size = DocxPt(8.5)
        run.font.color.rgb = RGBColor(120, 120, 112)


def _add_risk_cards(doc: Document, cards: list[Any]) -> None:
    cols = 2
    rows = (len(cards) + cols - 1) // cols
    table = doc.add_table(rows=rows, cols=cols)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    for index, card in enumerate(cards):
        if not isinstance(card, dict):
            card = {"title": "Risk", "body": str(card)}
        cell = table.rows[index // cols].cells[index % cols]
        _shade_cell(cell, "FFFFFF")
        _set_cell_border(cell, top={"val": "single", "sz": "4", "color": "E7E2D8"}, bottom={"val": "single", "sz": "4", "color": "E7E2D8"}, left={"val": "single", "sz": "4", "color": "E7E2D8"}, right={"val": "single", "sz": "4", "color": "E7E2D8"})
        p = cell.paragraphs[0]
        title_run = p.add_run(str(card.get("title", "Risk")))
        title_run.font.bold = True
        title_run.font.color.rgb = RGBColor(24, 58, 99)
        p.add_run("\n")
        body_run = p.add_run(str(card.get("body", "")))
        body_run.font.size = DocxPt(9.5)
    doc.add_paragraph()


def _add_footer_note(doc: Document, disclaimer: str, footer: str) -> None:
    table = doc.add_table(rows=1, cols=2)
    left, right = table.rows[0].cells
    _clear_cell(left)
    _clear_cell(right)
    left_run = left.paragraphs[0].add_run(disclaimer)
    left_run.font.size = DocxPt(8.5)
    left_run.font.color.rgb = RGBColor(120, 120, 112)
    right.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.RIGHT
    right_run = right.paragraphs[0].add_run(footer)
    right_run.font.size = DocxPt(8.5)
    right_run.font.color.rgb = RGBColor(120, 120, 112)


def _clear_cell(cell: Any) -> None:
    cell.text = ""


def _shade_cell(cell: Any, fill: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shading = tc_pr.find(qn("w:shd"))
    if shading is None:
        shading = OxmlElement("w:shd")
        tc_pr.append(shading)
    shading.set(qn("w:fill"), fill)


def _set_cell_border(cell: Any, **borders: dict[str, str]) -> None:
    tc = cell._tc
    tc_pr = tc.get_or_add_tcPr()
    tc_borders = tc_pr.first_child_found_in("w:tcBorders")
    if tc_borders is None:
        tc_borders = OxmlElement("w:tcBorders")
        tc_pr.append(tc_borders)
    for edge in ("top", "left", "bottom", "right", "insideH", "insideV"):
        edge_data = borders.get(edge)
        if edge_data is None:
            continue
        tag = f"w:{edge}"
        element = tc_borders.find(qn(tag))
        if element is None:
            element = OxmlElement(tag)
            tc_borders.append(element)
        for key, value in edge_data.items():
            element.set(qn(f"w:{key}"), str(value))


PPT_WARM_BG = "F7F5EF"
PPT_INK = "171717"
PPT_GRAY = "5E5C56"
PPT_MUTED = "747169"
PPT_NAVY = "183A63"
PPT_PALE_BLUE = "EEF4FB"
PPT_RULE = "E7E2D8"
PPT_SERIF = "Georgia"
PPT_SANS = "Aptos"
PPT_MONO = "Courier New"


def _create_warm_editorial_deck(title: str, slides: list[dict[str, Any]], target: Any) -> None:
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)
    prs.core_properties.author = "Code Assist"
    prs.core_properties.title = title
    prs.core_properties.subject = "Warm editorial presentation"

    deck_label = _deck_label(slides)
    _ppt_add_cover(prs, title=title, slides=slides, deck_label=deck_label)
    for index, item in enumerate(slides, start=2):
        _ppt_add_content_slide(prs, item=item, slide_number=index, deck_label=deck_label)
    prs.save(target)


def _deck_label(slides: list[dict[str, Any]]) -> str:
    for item in slides:
        label = item.get("deck_label") or item.get("footer")
        if label:
            return str(label).upper()
    return "AGENT ENGINEERING"


def _ppt_blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = _ppt_color(PPT_WARM_BG)
    return slide


def _ppt_add_cover(prs: Presentation, title: str, slides: list[dict[str, Any]], deck_label: str) -> None:
    slide = _ppt_blank_slide(prs)
    first = slides[0] if slides else {}
    kicker = str(first.get("cover_kicker") or first.get("kicker") or "KEYNOTE · 2026").upper()
    subtitle = str(first.get("subtitle") or first.get("body") or "")
    meta = str(first.get("meta") or f"code assist deck · A4 landscape · {len(slides) + 1} slides")

    _ppt_text(slide, kicker, 1.0, 1.35, 4.5, 0.25, size=12, font=PPT_MONO, color=PPT_NAVY, spacing=5.5)
    _ppt_text(slide, title, 1.0, 1.85, 6.4, 1.55, size=42, font=PPT_SERIF, color=PPT_INK, bold=False)
    if subtitle:
        _ppt_text(slide, subtitle, 1.0, 3.25, 5.7, 0.65, size=18, font=PPT_SERIF, color=PPT_GRAY)
    _ppt_line(slide, 1.0, 4.35, 1.0, 0, color=PPT_NAVY, width=2.0)
    _ppt_text(slide, meta, 1.0, 4.62, 5.6, 0.28, size=11, font=PPT_SERIF, color=PPT_MUTED)

    _ppt_cover_rings(slide)
    _ppt_footer(slide, slide_number=1, deck_label=deck_label)


def _ppt_add_content_slide(prs: Presentation, item: dict[str, Any], slide_number: int, deck_label: str) -> None:
    slide = _ppt_blank_slide(prs)
    section = str(item.get("section") or item.get("kicker") or f"{slide_number - 1:02d} · SECTION").upper()
    title = str(item.get("title") or "Slide")
    subtitle = str(item.get("subtitle") or item.get("body") or item.get("lead") or "")
    bullets = _as_list(item.get("bullets", []))

    _ppt_text(slide, section, 1.0, 0.72, 4.8, 0.3, size=11, font=PPT_MONO, color=PPT_MUTED, spacing=5.5)
    _ppt_text(slide, title, 1.0, 1.18, 6.8, 1.25, size=39, font=PPT_SERIF, color=PPT_INK)

    has_visual = any(item.get(key) for key in ("diagram", "chart", "code"))
    if has_visual:
        if subtitle:
            _ppt_text(slide, subtitle, 1.0, 2.86, 5.45, 0.75, size=20, font=PPT_SERIF, color=PPT_GRAY)
        _ppt_add_numbered_bullets(slide, bullets, 1.0, 3.72, 5.2, 1.65)
        if item.get("quote"):
            _ppt_add_callout(slide, str(item["quote"]), 1.0, 5.35, 4.9, 0.85)
        if item.get("diagram"):
            _ppt_add_loop_diagram(slide, item.get("diagram"))
        elif item.get("chart"):
            _ppt_add_bar_chart(slide, item.get("chart"))
        elif item.get("code"):
            _ppt_add_code_block(slide, str(item["code"]), 7.0, 2.55, 5.25, 2.85)
    elif item.get("stats"):
        if subtitle:
            _ppt_text(slide, subtitle, 1.0, 2.9, 7.4, 0.92, size=21, font=PPT_SERIF, color=PPT_GRAY)
        _ppt_add_numbered_bullets(slide, bullets, 1.0, 3.95, 7.2, 1.35)
        _ppt_add_stats(slide, item.get("stats", []))
    elif item.get("quote"):
        _ppt_text(slide, str(item["quote"]), 3.55, 2.2, 6.3, 2.1, size=41, font=PPT_SERIF, color=PPT_INK, align=PP_ALIGN.CENTER)
        _ppt_line(slide, 6.2, 4.72, 0.95, 0, color=PPT_NAVY, width=1.6)
        if subtitle:
            _ppt_text(slide, subtitle, 3.55, 5.12, 6.3, 0.85, size=17, font=PPT_SERIF, color=PPT_GRAY, align=PP_ALIGN.CENTER)
    else:
        if subtitle:
            _ppt_text(slide, subtitle, 1.0, 2.85, 7.4, 1.0, size=21, font=PPT_SERIF, color=PPT_GRAY)
        _ppt_add_numbered_bullets(slide, bullets, 1.0, 4.02, 7.0, 1.85)
        if item.get("callout"):
            _ppt_add_callout(slide, str(item["callout"]), 1.0, 5.45, 5.2, 0.8)

    _ppt_footer(slide, slide_number=slide_number, deck_label=deck_label)


def _ppt_cover_rings(slide: Any) -> None:
    cx, cy = 9.8, 2.85
    _ppt_oval(slide, cx - 1.35, cy - 1.35, 2.7, 2.7, fill=None, line=PPT_MUTED, width=1.1)
    _ppt_oval(slide, cx - 0.95, cy - 0.95, 1.9, 1.9, fill=None, line=PPT_RULE, width=0.8)
    _ppt_oval(slide, cx - 0.58, cy - 0.58, 1.16, 1.16, fill=PPT_PALE_BLUE, line=PPT_NAVY, width=1.2)
    _ppt_line(slide, cx, cy - 1.35, 1.35, 0, color=PPT_NAVY, width=2.0)
    _ppt_text(slide, "Agent\nLOOP", cx - 0.38, cy - 0.23, 0.76, 0.48, size=9, font=PPT_SERIF, color=PPT_INK, align=PP_ALIGN.CENTER)


def _ppt_add_loop_diagram(slide: Any, diagram: Any) -> None:
    labels = ["PLAN", "ACT", "OBSERVE", "REFLECT"]
    if isinstance(diagram, dict):
        labels = [str(x) for x in diagram.get("labels", labels)][:4]
    cx, cy = 9.75, 4.0
    _ppt_oval(slide, cx - 1.55, cy - 1.55, 3.1, 3.1, fill=None, line=PPT_MUTED, width=1.0)
    _ppt_oval(slide, cx - 1.15, cy - 1.15, 2.3, 2.3, fill=None, line=PPT_NAVY, width=0.7, dash=True)
    positions = [(cx - 0.75, cy - 2.0), (cx + 1.3, cy - 0.25), (cx - 0.75, cy + 1.48), (cx - 2.8, cy - 0.25)]
    for label, (x, y) in zip(labels, positions):
        _ppt_box(slide, label, x, y, 1.5, 0.48)
    _ppt_text(slide, "CORE", cx - 0.45, cy - 0.12, 0.9, 0.24, size=9, font=PPT_MONO, color=PPT_NAVY, align=PP_ALIGN.CENTER)


def _ppt_add_bar_chart(slide: Any, chart: Any) -> None:
    if not isinstance(chart, dict):
        chart = {}
    title = str(chart.get("title", "LOADING STRATEGY vs ACCURACY")).upper()
    values = chart.get("values", [53, 85])
    labels = chart.get("labels", ["Flat\nloading", "Layered\nloading"])
    values = [float(v) for v in values[:3]]
    labels = [str(v) for v in labels[: len(values)]]
    x, y, w, h = 7.1, 2.55, 5.0, 3.0
    _ppt_text(slide, title, x + 0.65, y - 0.3, 4.2, 0.25, size=10, font=PPT_MONO, color=PPT_MUTED, spacing=4.2, align=PP_ALIGN.CENTER)
    _ppt_line(slide, x, y + h, w, 0, color=PPT_GRAY, width=1.0)
    _ppt_line(slide, x, y, 0, h, color=PPT_GRAY, width=1.0)
    max_v = max(max(values), 1)
    bar_w = 1.25
    gap = 0.65
    for index, value in enumerate(values):
        bx = x + 0.45 + index * (bar_w + gap)
        bh = (value / max_v) * (h - 0.42)
        by = y + h - bh
        fill = PPT_PALE_BLUE if index == len(values) - 1 else "ECE9E0"
        line = PPT_NAVY if index == len(values) - 1 else "ECE9E0"
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptxInches(bx), PptxInches(by), PptxInches(bar_w), PptxInches(bh))
        rect.fill.solid()
        rect.fill.fore_color.rgb = _ppt_color(fill)
        rect.line.color.rgb = _ppt_color(line)
        rect.line.width = PptxPt(1.2 if index == len(values) - 1 else 0.2)
        _ppt_text(slide, f"{value:g}%", bx, by - 0.3, bar_w, 0.25, size=10, font=PPT_SANS, color=PPT_NAVY if index == len(values) - 1 else PPT_MUTED, align=PP_ALIGN.CENTER)
        _ppt_text(slide, labels[index], bx, y + h + 0.2, bar_w, 0.55, size=10, font=PPT_MONO, color=PPT_NAVY if index == len(values) - 1 else PPT_MUTED, align=PP_ALIGN.CENTER)
    target = chart.get("target")
    if target is not None:
        target_v = float(target)
        ty = y + h - (target_v / max_v) * (h - 0.42)
        _ppt_line(slide, x, ty, w, 0, color=PPT_NAVY, width=0.8, dash=True)


def _ppt_add_code_block(slide: Any, code: str, x: float, y: float, w: float, h: float) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptxInches(x), PptxInches(y), PptxInches(w), PptxInches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _ppt_color("FBFAF6")
    shape.line.color.rgb = _ppt_color(PPT_RULE)
    shape.line.width = PptxPt(0.8)
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.margin_left = PptxInches(0.28)
    text_frame.margin_right = PptxInches(0.22)
    text_frame.margin_top = PptxInches(0.22)
    lines = code.strip().splitlines()[:11]
    for index, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = line
        paragraph.font.name = PPT_MONO
        paragraph.font.size = PptxPt(11)
        paragraph.font.color.rgb = _ppt_color(PPT_NAVY if any(token in line for token in ("function", "if", "return")) else PPT_GRAY)


def _ppt_add_stats(slide: Any, stats: Any) -> None:
    if not isinstance(stats, list):
        return
    x_positions = [1.0, 5.05, 9.1]
    for index, stat in enumerate(stats[:3]):
        if not isinstance(stat, dict):
            stat = {"value": str(stat), "label": ""}
        x = x_positions[index]
        _ppt_text(slide, str(stat.get("value", "")), x, 5.82, 2.9, 0.55, size=34, font=PPT_SERIF, color=PPT_NAVY)
        _ppt_text(slide, str(stat.get("label", "")), x, 6.35, 3.1, 0.35, size=11, font=PPT_SERIF, color=PPT_GRAY)


def _ppt_add_numbered_bullets(slide: Any, bullets: list[str], x: float, y: float, w: float, h: float) -> None:
    if not bullets:
        return
    text_frame = slide.shapes.add_textbox(PptxInches(x), PptxInches(y), PptxInches(w), PptxInches(h)).text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    for index, bullet in enumerate(bullets[:5], start=1):
        paragraph = text_frame.paragraphs[0] if index == 1 else text_frame.add_paragraph()
        paragraph.text = f"{index}.  {bullet}"
        paragraph.font.name = PPT_SERIF
        paragraph.font.size = PptxPt(15)
        paragraph.font.color.rgb = _ppt_color(PPT_GRAY)
        paragraph.space_after = PptxPt(9)
        if paragraph.runs:
            paragraph.runs[0].font.color.rgb = _ppt_color(PPT_NAVY)


def _ppt_add_callout(slide: Any, text: str, x: float, y: float, w: float, h: float) -> None:
    _ppt_line(slide, x, y, 0, h, color=PPT_NAVY, width=2.0)
    _ppt_text(slide, text, x + 0.3, y + 0.08, w - 0.3, h - 0.05, size=16, font=PPT_SERIF, color=PPT_GRAY)


def _ppt_footer(slide: Any, slide_number: int, deck_label: str) -> None:
    _ppt_text(slide, deck_label, 1.0, 6.86, 3.5, 0.22, size=10, font=PPT_MONO, color=PPT_MUTED, spacing=4.0)
    _ppt_text(slide, f"{slide_number:02d}", 12.25, 6.86, 0.55, 0.22, size=10, font=PPT_MONO, color=PPT_MUTED, align=PP_ALIGN.RIGHT)


def _ppt_box(slide: Any, text: str, x: float, y: float, w: float, h: float) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptxInches(x), PptxInches(y), PptxInches(w), PptxInches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _ppt_color("FFFFFF")
    shape.line.color.rgb = _ppt_color(PPT_GRAY)
    shape.line.width = PptxPt(0.7)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = PPT_SERIF
    p.font.size = PptxPt(9)
    p.font.color.rgb = _ppt_color(PPT_INK)


def _ppt_oval(slide: Any, x: float, y: float, w: float, h: float, fill: str | None, line: str, width: float, dash: bool = False) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, PptxInches(x), PptxInches(y), PptxInches(w), PptxInches(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _ppt_color(fill)
    else:
        shape.fill.background()
    shape.line.color.rgb = _ppt_color(line)
    shape.line.width = PptxPt(width)
    if dash:
        shape.line.dash_style = 4


def _ppt_line(slide: Any, x: float, y: float, w: float, h: float, color: str, width: float, dash: bool = False) -> None:
    line = slide.shapes.add_connector(1, PptxInches(x), PptxInches(y), PptxInches(x + w), PptxInches(y + h))
    line.line.color.rgb = _ppt_color(color)
    line.line.width = PptxPt(width)
    if dash:
        line.line.dash_style = 4


def _ppt_text(
    slide: Any,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float,
    font: str,
    color: str,
    bold: bool = False,
    align: Any = PP_ALIGN.LEFT,
    spacing: float | None = None,
) -> None:
    box = slide.shapes.add_textbox(PptxInches(x), PptxInches(y), PptxInches(w), PptxInches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = font
    p.font.size = PptxPt(size)
    p.font.bold = bold
    p.font.color.rgb = _ppt_color(color)
    if spacing is not None:
        for run in p.runs:
            run.font._element.set("spc", str(int(spacing * 100)))


def _ppt_color(hex_value: str) -> PptxRGBColor:
    value = hex_value.strip("#")
    return PptxRGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _parse_markdown_slides(markdown: str) -> tuple[str, list[dict[str, Any]]]:
    title = "Presentation"
    slides: list[dict[str, Any]] = []
    current: dict[str, Any] | None = None

    for raw_line in markdown.splitlines():
        line = raw_line.strip()
        if not line:
            continue
        if line.startswith("# ") and title == "Presentation":
            title = line[2:].strip() or title
            continue
        if line.startswith("## "):
            if current:
                slides.append(current)
            current = {"title": line[3:].strip() or "Slide", "bullets": []}
            continue
        if current is None:
            continue
        if line.startswith(("- ", "* ")):
            current["bullets"].append(line[2:].strip())
        elif len(current["bullets"]) < 5:
            current["bullets"].append(line)

    if current:
        slides.append(current)
    if not slides:
        slides.append({"title": "Summary", "bullets": [line for line in markdown.splitlines() if line.strip()][:5]})
    return title, slides
